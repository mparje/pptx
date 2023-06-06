import streamlit as st
import os
from pptx import Presentation
from pptx.util import Inches
import openai
import fitz

# Configurar la clave de la API de OpenAI
api_key = st.sidebar.text_input("Enter your OpenAI API key", type="password")

if not api_key:
    st.warning("Please enter a valid API key to continue..")
else:
    openai.api_key = api_key
    # Continuar con el resto del código que utiliza la clave de API
    
    
# Generar el resumen utilizando GPT-3
def generate_summary(pdf_path):
    # Leer el contenido del PDF y almacenarlo en una variable llamada 'content'
    content = read_pdf_content(pdf_path)

    # Llamar a la API de GPT-3 para generar el resumen
    response = openai.Completion.create(
        engine='text-davinci-003',
        prompt=content,
        max_tokens=150,
        temperature=0.5,
        top_p=1.0,
        frequency_penalty=0.0,
        presence_penalty=0.0
    )

    summary = response.choices[0].text.strip()
    return summary


# Leer el contenido del PDF
def read_pdf_content(pdf_path):
    # Implementa la lógica para leer el contenido del PDF y devolverlo como una cadena de texto
    # Puedes utilizar una biblioteca como PyPDF2 o pdfminer.six
    pass


# Convertir el PDF a imágenes utilizando PyMuPDF
def convert_pdf_to_images(pdf_path):
    doc = fitz.open(pdf_path)
    image_paths = []

    for i in range(len(doc)):
        page = doc.load_page(i)
        pix = page.get_pixmap()
        image_path = f'page_{i}.png'
        pix.save(image_path)
        image_paths.append(image_path)

    return image_paths


# Generar el archivo PPTX utilizando el resumen y las imágenes
def generate_pptx(summary, images, pptx_path):
    presentation = Presentation()

    # Agregar el resumen como el contenido de la primera diapositiva
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = 'Resumen del PDF'
    slide.placeholders[1].text = summary

    # Agregar las imágenes a las diapositivas
    for image_path in images:
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.add_picture(image_path, Inches(1), Inches(1), width=Inches(8), height=Inches(6))

    # Guardar la presentación como un archivo PPTX
    presentation.save(pptx_path)


# Aplicación de Streamlit
def main():
    st.title("Generador de Presentaciones PPTX desde PDF")
    pdf_file = st.file_uploader("Cargar archivo PDF", type=["pdf"])

    if pdf_file is not None:
        # Guardar el archivo PDF en el sistema de archivos
        with open('temp.pdf', 'wb') as f:
            f.write(pdf_file.getvalue())

        # Generar el resumen del PDF
        summary = generate_summary('temp.pdf')

        # Convertir las páginas del PDF a imágenes
        images = convert_pdf_to_images('temp.pdf')

        # Generar la presentación PPTX utilizando el resumen y las imágenes
        generate_pptx(summary, images, 'output.pptx')

        # Descargar el archivo PPTX
        with open('output.pptx', 'rb') as f:
            st.download_button('Descargar PPTX', f.read(), file_name='output.pptx')

        # Eliminar los archivos temporales
        os.remove('temp.pdf')
        for image_path in images:
            os.remove(image_path)


if __name__ == '__main__':
    main()
